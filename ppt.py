from pptx import Presentation
# Correcting the issue with the architecture slide and regenerating the presentation

# Create a new presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Image-Based Entity Extraction for Digital Marketplaces"
subtitle.text = "A Machine Learning Solution for Accurate Data Extraction"

# Slide 2: Project Overview
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Project Overview"
content.text = (
    "Extracts key entity values (e.g., weight, dimensions) from product images.\n"
    "Crucial for fields like e-commerce, healthcare, and content moderation.\n"
    "Addresses the challenge of missing textual descriptions in digital marketplaces."
)

# Slide 3: Problem Statement
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Problem Statement"
content.text = (
    "Many digital products lack detailed textual descriptions.\n"
    "Manual extraction is time-consuming and error-prone.\n"
    "Need for an automated solution to extract precise information from images."
)

# Slide 4: Proposed Solution
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Proposed Solution"
content.text = (
    "Machine Learning model to extract entities directly from images.\n"
    "Preprocessing module for image enhancement.\n"
    "Data formatting and validation to ensure accuracy.\n"
    "Scalable deployment using Vultr cloud services."
)

# Slide 5: Architecture Diagram
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "System Architecture"
# Adding placeholder for architecture diagram
txBox = slide.shapes.add_textbox(left=300, top=200, width=5000, height=1000)
tf = txBox.text_frame
tf.text = "Insert Architecture Diagram Here"

# Slide 6: Key Modules
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Key Modules"
content.text = (
    "• User Interface (Web/API) for image upload.\n"
    "• Image Preprocessing for standardization.\n"
    "• ML Model for entity extraction.\n"
    "• Entity Formatting Module for data validation.\n"
    "• Vultr Compute Instances and Object Storage."
)

# Slide 7: Use of Vultr Services
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Use of Vultr Services"
content.text = (
    "• Compute Instances for scalable model inference.\n"
    "• Object Storage for raw images and extracted data.\n"
    "• Ensures high availability and performance."
)

# Slide 8: Data Flow
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Data Flow"
content.text = (
    "1. User uploads an image via UI/API.\n"
    "2. Image is preprocessed and sent to the ML model.\n"
    "3. Entities are extracted and formatted.\n"
    "4. Data is stored in Vultr Object Storage and database."
)

# Slide 9: Target Audience
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Target Audience"
content.text = (
    "• E-commerce platforms for detailed product descriptions.\n"
    "• Healthcare for accurate product data extraction.\n"
    "• Content moderation teams for automated verification."
)

# Slide 10: Expected Outcomes
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Expected Outcomes"
content.text = (
    "• Accurate and automated entity extraction from images.\n"
    "• Enhanced efficiency in data processing.\n"
    "• Scalable solution for various industries."
)

# Slide 11: Conclusion
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "• A reliable ML-based solution for extracting essential data from images.\n"
    "• Scalability and efficiency ensured with Vultr services.\n"
    "• Significant impact on e-commerce, healthcare, and beyond."
)

# Save the presentation
pptx_path = "Project_Presentation.pptx"
presentation.save(pptx_path)
pptx_path
